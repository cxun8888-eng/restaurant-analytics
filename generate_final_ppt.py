"""生成项目展示PPT — 淡色调 + 截图占位"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== 淡色配色 =====
C_BG = RGBColor(0xFA, 0xFA, 0xFA)           # 极浅灰底
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # 白
C_TEXT = RGBColor(0x1a, 0x1a, 0x2e)           # 深黑文字
C_GRAY = RGBColor(0x6B, 0x72, 0x80)           # 灰色
C_LIGHT = RGBColor(0xE5, 0xE7, 0xEB)          # 浅灰线
C_BLUE = RGBColor(0x60, 0xA5, 0xFA)           # 柔和蓝
C_GREEN = RGBColor(0x34, 0xD3, 0x99)          # 柔和绿
C_ORANGE = RGBColor(0xFB, 0xBF, 0x24)         # 柔和橙
C_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)         # 柔和紫
C_PINK = RGBColor(0xFB, 0x92, 0x3C)           # 柔和粉
C_CARD_BG = RGBColor(0xF8, 0xF9, 0xFA)        # 卡片浅灰

ACCENTS = [C_BLUE, C_GREEN, C_ORANGE, C_PURPLE, C_PINK]


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG


def tb(slide, left, top, width, height, text, size=14, color=C_TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def screenshot_placeholder(slide, x, y, w, h, label="截图位置"):
    """虚线框截图占位"""
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)  # rounded rect
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_WHITE
    shape.line.color.rgb = C_LIGHT
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 2  # dash
    # 文字
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[ {label} ]"
    p.font.size = Pt(12)
    p.font.color.rgb = C_GRAY
    p.alignment = PP_ALIGN.CENTER


def card(slide, x, y, w, h, title, lines, accent=C_BLUE):
    """浅色卡片"""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_WHITE
    shape.line.color.rgb = C_LIGHT
    shape.line.width = Pt(1)
    # 顶部色条
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    # 内容
    tb(slide, x + 0.2, y + 0.2, w - 0.4, 0.3, title, size=14, color=C_TEXT, bold=True)
    tb(slide, x + 0.2, y + 0.6, w - 0.4, h - 0.8, lines, size=10, color=C_GRAY)


def section_title(slide, text, y=0.3):
    tb(slide, 0.5, y, 12, 0.6, text, size=28, color=C_TEXT, bold=True)


def algorithm_tag(slide, x, y, text, accent=C_BLUE):
    """算法标签"""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(len(text)*0.15+0.3), Inches(0.3))
    shape.fill.solid(); shape.fill.fore_color.rgb = accent; shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(8); p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.CENTER


# ============================================================
# Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
shape = slide.shapes.add_shape(1, Inches(2), Inches(2.5), Inches(9.3), Inches(0.03))
shape.fill.solid(); shape.fill.fore_color.rgb = C_BLUE; shape.line.fill.background()

tb(slide, 2, 1.2, 9.3, 0.8, '餐饮多平台经营数据分析系统', size=38, color=C_TEXT, bold=True, align=PP_ALIGN.CENTER)
tb(slide, 2, 2.8, 9.3, 0.5, '从数据到决策：全链路经营分析平台', size=16, color=C_GRAY, align=PP_ALIGN.CENTER)
tb(slide, 2, 4.5, 9.3, 0.4, '数据科学与大数据技术 · 个人项目', size=14, color=C_BLUE, align=PP_ALIGN.CENTER)
tb(slide, 2, 5.2, 9.3, 0.3, 'Python | Streamlit | Pandas | Scikit-learn | Plotly', size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 2: 项目概述 + 截图占位
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, '项目概述')

tb(slide, 0.5, 1.2, 7, 1.8,
   '商家从美团、饿了么、微信点单导出订单数据，上传到本平台',
   size=13, color=C_GRAY)
tb(slide, 0.5, 1.6, 7, 1.8,
   '自动完成：数据清洗 → 多维分析 → 算法建模 → 可视化 → 智能经营建议',
   size=13, color=C_TEXT, bold=True)

# 6模块
mods = ['数据上传', '经营概览', '商品分析', '用户分析', '智能预测', '分析报告']
for i, m in enumerate(mods):
    x = 0.5 + i * 1.2
    tag_color = ACCENTS[i % 5]
    shape = slide.shapes.add_shape(1, Inches(x), Inches(2.2), Inches(1.1), Inches(0.35))
    shape.fill.solid(); shape.fill.fore_color.rgb = tag_color; shape.line.fill.background()
    tf = shape.text_frame; p = tf.paragraphs[0]; p.text = m; p.font.size = Pt(9); p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.CENTER

# 截图占位
screenshot_placeholder(slide, 8, 1.0, 5, 3.2, '项目首页截图')

tb(slide, 0.5, 4.8, 7, 1.5, '支持多平台数据统一分析 | 6种算法模型 | 已部署公网可访问', size=11, color=C_BLUE)

# ============================================================
# Slide 3: 数据上传
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, '📤 数据上传 — ETL 数据管道')

# 左侧：流程
steps = [
    ('01', '文件加载', 'CSV/Excel 自动检测编码'),
    ('02', '列名标准化', '多平台表头自动映射'),
    ('03', '数据校验', '必填字段 + 日期 + 非负检查'),
    ('04', '数据清洗', '去重 · 缺失值填充 · 字段推算'),
    ('05', '异常检测', 'IQR 方法标记异常订单'),
]
for i, (num, title, desc) in enumerate(steps):
    y = 1.3 + i * 0.85
    tb(slide, 0.8, y, 0.4, 0.3, num, size=22, color=C_BLUE, bold=True)
    tb(slide, 1.3, y, 1.8, 0.25, title, size=13, color=C_TEXT, bold=True)
    tb(slide, 1.3, y + 0.3, 3, 0.3, desc, size=10, color=C_GRAY)

# 右侧：IQR 说明
tb(slide, 5.5, 1.3, 3, 0.3, 'IQR 异常检测', size=15, color=C_TEXT, bold=True)
tb(slide, 5.5, 1.8, 3, 2.5,
   '''Q1 = 第 25 百分位
Q3 = 第 75 百分位
IQR = Q3 - Q1

正常范围：
Q1 - 1.5 × IQR  ~  Q3 + 1.5 × IQR

1.5 是经验值
等效正态分布下 2.7σ''', size=11, color=C_GRAY)

# 截图占位
screenshot_placeholder(slide, 9, 1.0, 3.8, 3.5, '数据上传页面截图')

algorithm_tag(slide, 0.5, 6.0, 'IQR 四分位距', C_BLUE)
algorithm_tag(slide, 1.8, 6.0, 'ETL 管道', C_BLUE)

# ============================================================
# Slide 4: 经营概览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, '📊 经营概览 — 描述性统计分析')

features = [
    ('核心指标', '总营收 · 订单数\n客单价 · 退款率\n今日/昨日对比'),
    ('营收趋势', '日营收折线图\n7日移动平均\n日环比柱状图'),
    ('时段热力图', '工作日 × 小时\n交叉分析\n识别高峰时段'),
    ('平台对比', '美团 vs 微信 vs 饿了么\n客单价/退款率对比\n营收占比'),
    ('统计五数', '均值 · 中位数 · 标准差\n偏度 · 峰度\n分布形态分析'),
]
for i, (title, desc) in enumerate(features):
    card(slide, 0.3 + i * 2.6, 1.3, 2.4, 2.8, title, desc, ACCENTS[i])

# 截图占位
screenshot_placeholder(slide, 1, 4.5, 11.3, 2.0, '经营概览页面截图')

algorithm_tag(slide, 0.5, 6.8, '描述性统计', C_BLUE)
algorithm_tag(slide, 1.8, 6.8, '时间序列分析', C_GREEN)

# ============================================================
# Slide 5: 商品分析
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, '🔗 商品分析 — Apriori 关联规则')

# 左侧：原理
tb(slide, 0.5, 1.3, 5.5, 0.3, '算法原理', size=15, color=C_TEXT, bold=True)
tb(slide, 0.8, 1.8, 5.5, 2,
   '''1. 构建 订单 × 商品 one-hot 矩阵
2. 扫描频繁项集（先验剪枝，避免穷举）
3. 计算关联规则三个指标
4. 按 Lift 排序，自动生成套餐建议''', size=11, color=C_GRAY)

# 三个指标卡片
metrics = [
    ('支持度', 'Support', 'P(A∩B)', 'A+B 同时出现 ÷ 总订单\n衡量"普遍不普遍"'),
    ('置信度', 'Confidence', 'P(B|A)', '买了A也买B的比例\n衡量"靠不靠谱"'),
    ('提升度', 'Lift', 'P(B|A)/P(B)', '排除了高频商品干扰\n>1 = 正相关，越大越强'),
]
for i, (cn, en, formula, desc) in enumerate(metrics):
    x = 7 + i * 2.2
    card(slide, x, 1.3, 2.0, 2.5, f'{cn} ({en})', f'{formula}\n\n{desc}', ACCENTS[i])

# 截图占位
screenshot_placeholder(slide, 0.5, 4.2, 6, 2.5, '商品销量排行截图')
screenshot_placeholder(slide, 7, 4.2, 5.8, 2.5, '关联规则结果截图')

algorithm_tag(slide, 0.5, 6.9, 'Apriori', C_GREEN)
algorithm_tag(slide, 1.3, 6.9, 'Support/Confidence/Lift', C_GREEN)

# ============================================================
# Slide 6: 用户分析
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, '👤 用户分析 — RFM 分层 + K-Means 聚类')

# 左侧RFM
tb(slide, 0.5, 1.3, 5.5, 0.3, 'RFM 用户分层', size=15, color=C_TEXT, bold=True)
tb(slide, 0.8, 1.8, 5.5, 2.5,
   '''R (Recency)  = 最近一次消费距今多少天
F (Frequency) = 累计消费了多少次
M (Monetary)  = 累计消费了多少钱

→ 分位数三等分，每维打 1~3 分
→ 3×3×3 = 27 种组合 → 归纳为 8 类用户
→ 每类用户匹配不同运营策略''', size=11, color=C_GRAY)

# 8类表格
seg_data = [
    ('R高F高M高', '重要价值', 'VIP 待遇'),
    ('R高F高M中', '潜力客户', '推高价新品'),
    ('R中F低M高', '重要保持', '发券召回'),
    ('R低F高M高', '重要挽留', '大力度折扣'),
    ('R低F低M低', '流失客户', '低成本触达'),
]
tbl = slide.shapes.add_table(len(seg_data)+1, 3, Inches(0.8), Inches(4.6), Inches(5.5), Inches(0.25*6))
headers = ['RFM 组合', '分层标签', '运营策略']
for c, h in enumerate(headers):
    cell = tbl.table.cell(0, c); cell.text = h
    for p in cell.text_frame.paragraphs: p.font.size = Pt(9); p.font.bold = True
for r, row in enumerate(seg_data):
    for c, val in enumerate(row):
        cell = tbl.table.cell(r+1, c); cell.text = val
        for p in cell.text_frame.paragraphs: p.font.size = Pt(9); p.font.color.rgb = C_GRAY

# 右侧K-Means
tb(slide, 7, 1.3, 5.5, 0.3, 'K-Means 聚类验证', size=15, color=C_TEXT, bold=True)
tb(slide, 7, 1.8, 5.5, 2,
   '''肘部法则确定最优 K 值
无监督聚类自然分组
交叉表对比 RFM vs K-Means
→ 两种方法结果高度一致
→ RFM 分层客观可靠''', size=11, color=C_GRAY)

# 截图占位
screenshot_placeholder(slide, 7, 4.2, 5.5, 2.5, 'RFM 3D散点图截图')

algorithm_tag(slide, 0.5, 6.9, 'RFM 模型', C_BLUE)
algorithm_tag(slide, 1.5, 6.9, 'K-Means 聚类', C_PURPLE)

# ============================================================
# Slide 7: 预测 + 异常检测
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, '🔮 智能预测 — 随机森林 + 异常检测')

# 预测
tb(slide, 0.5, 1.3, 6, 0.3, '营收预测：随机森林回归', size=15, color=C_TEXT, bold=True)
tb(slide, 0.8, 1.8, 6, 2.5,
   '''思路：时间序列 → 监督学习

构造 17 维时间特征：
  趋势 · 周期 · 季节 · 滞后 · 滚动统计

训练随机森林 (100棵树)
递归预测未来 14 天
MAPE < 9%     95% 置信区间''', size=11, color=C_GRAY)

# 异常检测
tb(slide, 0.5, 4.5, 6, 0.3, '异常检测：Isolation Forest', size=15, color=C_TEXT, bold=True)
tb(slide, 0.8, 5.0, 6, 1.5,
   '''核心思想：随机切割 → 异常点容易被隔离

5 维检测特征：
  订单金额 · 商品种类数 · 商品总数量
  平均单价 · 折扣金额''', size=11, color=C_GRAY)

# 对比表
tb(slide, 7, 1.3, 5.5, 0.3, '算法选型对比', size=15, color=C_TEXT, bold=True)
tb(slide, 7, 1.8, 5.5, 3,
   '''LSTM：数据量不够 (90天)，可解释性差
Prophet：本质调包，不能体现建模能力
ARIMA：节假日处理弱，需人工哑变量

→ 随机森林 + 特征工程
  = 把时序转化为 ML 问题的通用方法论
  = 机器学习核心算法，面试官认可
  = 可解释：特征重要性排序''', size=11, color=C_GRAY)

# 截图占位
screenshot_placeholder(slide, 7, 5.0, 5.5, 2, '预测图表截图')

algorithm_tag(slide, 0.5, 6.9, '随机森林回归', C_ORANGE)
algorithm_tag(slide, 1.8, 6.9, 'Isolation Forest', C_PINK)
algorithm_tag(slide, 3.3, 6.9, '特征工程 (17维)', C_ORANGE)

# ============================================================
# Slide 8: 分析报告
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, '📋 分析报告 — 数据 → 洞察 → 决策')

report_items = [
    ('经营概览', '营收 · 订单 · 客单价\n统计五数', C_BLUE),
    ('异常提醒', '异常订单标记\n退款率/日环比预警', C_PINK),
    ('商品分析', '畅销排行 · 滞销关注\n套餐搭配建议', C_GREEN),
    ('用户洞察', '分层分布 · 风险用户\n新客引导建议', C_PURPLE),
    ('预测展望', '未来7天预测\n趋势判断 · 备货建议', C_ORANGE),
]
for i, (title, desc, accent) in enumerate(report_items):
    card(slide, 0.3 + i * 2.6, 1.3, 2.4, 2.2, title, desc, accent)

# 闭环
steps = ['上传 CSV', '清洗 ETL', '分析建模', '可视化', '诊断报告', '经营决策']
for i, s in enumerate(steps):
    x = 0.8 + i * 2.1
    tb(slide, x, 4.0, 1.8, 0.3, s, size=13, color=C_TEXT, bold=True, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        tb(slide, x + 1.8, 4.0, 0.3, 0.3, '→', size=16, color=C_GRAY, align=PP_ALIGN.CENTER)

tb(slide, 2, 4.8, 9, 0.5, '数据分析的最终价值不在于出图表，而在于输出可执行的决策', size=14, color=C_BLUE, bold=True, align=PP_ALIGN.CENTER)

# 截图占位
screenshot_placeholder(slide, 2, 5.5, 9.3, 1.5, '分析报告页面截图')

# ============================================================
# Slide 9: 算法总览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, '算法模型总览')

algos = [
    ('IQR', '四分位距', 'Q1-1.5IQR ~ Q3+1.5IQR', '数据管道初筛', C_BLUE),
    ('Apriori', '关联规则', 'Support / Confidence / Lift', '菜品搭配推荐', C_GREEN),
    ('RFM', '用户分层', '3维打分 → 8类分层', '差异化运营策略', C_BLUE),
    ('K-Means', '聚类', '肘部法则 · 无监督', '交叉验证RFM', C_PURPLE),
    ('随机森林', '集成回归', '100棵树 · 17维特征', '营收预测 MAPE<9%', C_ORANGE),
    ('Isolation\nForest', '异常检测', '多维空间切割', '异常订单识别', C_PINK),
]
for i, (name, cat, detail, usage, accent) in enumerate(algos):
    col = i % 3; row = i // 3
    x = 0.5 + col * 4.3; y = 1.3 + row * 2.9
    card(slide, x, y, 3.8, 2.5, f'{name} ({cat})', f'{detail}\n\n用途：{usage}', accent)

# ============================================================
# Slide 10: 技术栈
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, '技术栈 & 部署')

stacks = [
    ('Python', 'Streamlit · Pandas · NumPy\nScikit-learn · SciPy'),
    ('数据挖掘', 'mlxtend (Apriori)\n关联规则 · 频繁项集'),
    ('可视化', 'Plotly\n交互式图表'),
    ('部署', 'Hugging Face Spaces\nDocker · Git'),
    ('数据源', '模拟数据生成器\n90天 · 500用户 · 30+商品'),
]
for i, (title, desc) in enumerate(stacks):
    x = 0.5 + i * 2.6
    card(slide, x, 1.3, 2.4, 2.5, title, desc, ACCENTS[i])

screenshot_placeholder(slide, 2, 4.3, 9.3, 2.5, 'Hugging Face 部署页面截图 / 项目完整页面截图')

tb(slide, 2, 6.5, 9.3, 0.4, '数据科学与大数据技术 · 个人项目 · https://huggingface.co/spaces/momo815/restaurant-analytics',
   size=10, color=C_GRAY, align=PP_ALIGN.CENTER)

# ===== 保存 =====
out = os.path.join(os.path.expanduser('~/Desktop'), '餐饮数据分析_项目展示.pptx')
prs.save(out)
print(f'Saved: {out}')
