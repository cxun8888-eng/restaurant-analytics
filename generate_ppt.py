"""生成项目展示PPT — 餐饮多平台经营数据分析系统"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== 颜色方案 =====
C_DARK = RGBColor(0x1a, 0x1a, 0x2e)
C_PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
C_ACCENT = RGBColor(0x10, 0xB9, 0x81)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY = RGBColor(0x9c, 0xa3, 0xaf)
C_BG = RGBColor(0xf8, 0xf9, 0xfa)
C_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
C_RED = RGBColor(0xEF, 0x44, 0x44)


def add_bg(slide, color=C_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_shape(slide, left, top, width, height, color=C_PRIMARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, icon, title, desc, bg_color=None):
    if bg_color is None:
        bg_color = C_DARK if title else C_PRIMARY
    shape = add_shape(slide, left, top, width, height, bg_color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if icon:
        run = tf.paragraphs[0].add_run()
        run.text = icon
        run.font.size = Pt(28)
        run.font.color.rgb = C_WHITE
    if title:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = C_WHITE
    if desc:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(10)
        run.font.color.rgb = C_WHITE


# ============================================================
# Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)

# 装饰线
add_shape(slide, 2, 2.8, 9.3, 0.03, C_PRIMARY)

add_text_box(slide, 2, 1.8, 9.3, 0.8, '餐饮多平台经营数据分析系统', font_size=40, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, 2, 3.0, 9.3, 0.6, '从CSV上传到智能经营建议 — 全链路数据闭环', font_size=18, color=C_GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, 2, 4.0, 9.3, 0.5, '数据科学与大数据技术 · 个人项目', font_size=14, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 2, 5.2, 9.3, 0.4, 'Python | Streamlit | Pandas | Scikit-learn | Plotly', font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 2: 项目背景 & 目标
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_text_box(slide, 0.8, 0.5, 11, 0.6, '项目背景与目标', font_size=28, color=C_DARK, bold=True)

add_text_box(slide, 0.8, 1.5, 5.5, 2.5,
    '痛点\n\n'
    '商家从美团、饿了么、微信点单导出大量订单CSV，'
    '数据''死''在表格里，缺乏有效分析手段。\n\n'
    '→ 哪些菜好卖？哪些该下架？\n'
    '→ 顾客多久来一次？谁快流失了？\n'
    '→ 明天预计卖多少？该备多少货？', font_size=14, color=C_DARK)

add_text_box(slide, 7, 1.5, 5.5, 2.5,
    '解决方案\n\n'
    'Web端一键上传 → 自动清洗 → 多维分析 → 智能建议\n\n'
    '✓ 经营概览：营收趋势、时段分析、平台对比\n'
    '✓ 商品分析：销量排行、关联规则、套餐推荐\n'
    '✓ 用户分析：RFM分层、K-Means聚类\n'
    '✓ 智能预测：营收预测、异常订单检测\n'
    '✓ 分析报告：一键生成自然语言诊断建议', font_size=14, color=C_DARK)

# 底部技术标签
add_shape(slide, 0.8, 5.6, 11.7, 1.2, C_BG)
add_text_box(slide, 1.2, 5.8, 11, 1, '技术栈：Python | Streamlit | Pandas | NumPy | Scikit-learn | Plotly | mlxtend | SciPy\n'
           '数据来源：模拟数据（90天×日均120单，500位顾客，30+商品，3个平台） | 公网地址：已部署至 Streamlit Community Cloud',
           font_size=12, color=C_DARK)

# ============================================================
# Slide 3: 系统架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_text_box(slide, 0.8, 0.5, 11, 0.6, '系统架构与数据流', font_size=28, color=C_DARK, bold=True)

# 数据流: 箭头式
steps = ['上传\nCSV/Excel', 'ETL管道\n清洗校验', '特征工程\nRFM/时段', '算法建模\nML模型', '可视化\nPlotly', '诊断报告\n自然语言']
for i, step in enumerate(steps):
    x = 0.5 + i * 2.1
    add_card(slide, x, 1.6, 1.8, 1.2, '', step, '', C_PRIMARY if i < 3 else C_ACCENT if i < 5 else C_ORANGE)
    if i < len(steps) - 1:
        add_text_box(slide, x + 1.85, 2.0, 0.3, 0.3, '→', font_size=20, color=C_GRAY, align=PP_ALIGN.CENTER)

# 模块详情
modules = [
    ('数据上传', 'CSV/Excel导入\n列名自动识别\nIQR异常检测\n数据质量报告'),
    ('经营概览', '核心指标卡片\n营收趋势图\n时段热力图\n平台对比'),
    ('商品分析', '销量Top/Bottom\n品类营收占比\nApriori关联规则\n套餐搭配建议'),
    ('用户分析', 'RFM三维打分\n8类用户分层\nK-Means聚类验证\n差异化运营策略'),
    ('智能预测', '随机森林回归\n17维时间特征\nIsolation Forest\n多维异常检测'),
    ('分析报告', '一键生成报告\n自然语言建议\n异常提醒\n趋势判断'),
]
for i, (title, desc) in enumerate(modules):
    x = 0.5 + i * 2.1
    add_card(slide, x, 3.5, 1.8, 3.2, '', title, desc, C_DARK if i == 5 else C_PRIMARY if i == 0 else C_ACCENT)

# ============================================================
# Slide 4: 核心算法详解 - RFM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_text_box(slide, 0.8, 0.5, 6, 0.6, 'RFM 用户分层模型', font_size=28, color=C_DARK, bold=True)

add_text_box(slide, 0.8, 1.4, 5.5, 2.5,
    '三维度\n'
    'R (Recency)  = 最近一次消费距今多少天\n'
    'F (Frequency)= 累计消费了多少次\n'
    'M (Monetary) = 累计消费了多少钱\n\n'
    '分位数三等分 → 每维打1-3分\n'
    '3×3×3=27种组合 → 归纳为8类用户', font_size=14, color=C_DARK)

# 8类用户表格
table_data = [
    ('R分', 'F分', 'M分', '分层', '策略'),
    ('高(3)', '高(3)', '高(2-3)', '重要价值', 'VIP待遇'),
    ('高(3)', '中-高(2-3)', '中-高(2-3)', '潜力客户', '推高价新品'),
    ('中-高(2-3)', '低(1)', '中-高(2-3)', '重要保持', '发优惠券召回'),
    ('低(1)', '高(2-3)', '高(2-3)', '重要挽留', '大力度折扣'),
    ('低(1)', '低(1)', '低(1)', '流失客户', '低成本触达'),
]
rows = len(table_data)
cols = len(table_data[0])
table = slide.shapes.add_table(rows, cols, Inches(7), Inches(1.4), Inches(5.5), Inches(0.3 * rows)).table
for r, row in enumerate(table_data):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.CENTER

add_text_box(slide, 7, 3.8, 5.5, 1.2, '交叉验证：K-Means 聚类\n用肘部法则确定K值 → 无监督聚类自然分组\n→ 与RFM规则分层高度一致 → 分层客观可靠',
           font_size=13, color=C_DARK)

add_text_box(slide, 0.8, 4.5, 5.5, 2, '业务价值\n\n每类用户匹配不同运营策略。\n"重要挽留"—曾经高消费但很久没来→最高优先级召回。\n\n面试要点：RFM是CRM经典模型，\nAmazon/美团都在用，不需要训练数据。',
           font_size=13, color=C_DARK)

# ============================================================
# Slide 5: 核心算法详解 - Apriori
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_text_box(slide, 0.8, 0.5, 6, 0.6, 'Apriori 关联规则挖掘', font_size=28, color=C_DARK, bold=True)

add_text_box(slide, 0.8, 1.4, 5.5, 2.5,
    '目标：从海量订单中发现"买了A的顾客大概率也买B"\n\n'
    '三个核心指标\n\n'
    '支持度 (Support)  = A和B同时出现的订单 ÷ 总订单\n'
    '置信度 (Confidence)= 买了A的订单中也买了B的比例\n'
    '提升度 (Lift)       = 置信度 ÷ B的自然出现概率\n'
    '                          >1正相关 | >2强关联', font_size=14, color=C_DARK)

# 示例
add_shape(slide, 0.8, 4.2, 5.5, 1.5, C_BG)
add_text_box(slide, 1.2, 4.3, 4.8, 1.3,
    '示例：发现"酸辣土豆丝 + 可乐" Lift=3.2\n'
    '→ 点了土豆丝的人买可乐的概率是随机顾客的3.2倍\n'
    '→ 打包成套餐，定价略低于单点总价\n'
    '→ 提客单价 + 顾客觉得划算', font_size=13, color=C_DARK)

# 气泡图说明
add_text_box(slide, 7, 1.4, 5.5, 4,
    '算法流程\n\n'
    '1. 构建 订单×商品 one-hot矩阵\n'
    '   订单行，商品列，买了=1，没买=0\n\n'
    '2. 扫描频繁项集\n'
    '   利用先验性质剪枝\n'
    '   不频繁的项集的超集也一定不频繁\n\n'
    '3. 计算关联规则\n'
    '   对每个频繁项集计算Support/Confidence/Lift\n'
    '   按Lift降序排列\n\n'
    '4. 自动生成建议\n'
    '   Lift>2 → "建议打包为套餐"\n'
    '   Lift>3 → "强烈建议"', font_size=13, color=C_DARK)

add_text_box(slide, 7, 5.8, 5.5, 1, '面试要点：Lift是最关键的指标\nSupport只看频率，Confidence被高频商品误导\nLift排除了偏差，衡量真正相关性', font_size=12, color=C_PRIMARY)

# ============================================================
# Slide 6: 营收预测 & 异常检测
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_text_box(slide, 0.8, 0.5, 6, 0.6, '营收预测 & 异常订单检测', font_size=28, color=C_DARK, bold=True)

# 预测部分
add_text_box(slide, 0.8, 1.4, 5.5, 3,
    '营收预测 — 随机森林回归\n\n'
    '思路：时间序列 → 监督学习\n'
    '构造17维时间特征：\n'
    '• 趋势特征：day_num（第几天）\n'
    '• 周期特征：weekday one-hot（7维）\n'
    '• 季节特征：month、day_of_month\n'
    '• 滞后特征：lag_1/2/3/7（前N天营收）\n'
    '• 滚动统计：rolling_mean_7、rolling_std_7\n'
    '• 周末标记：is_weekend\n\n'
    '训练随机森林（100棵树，max_depth=5）\n'
    '→ 递归预测未来14天 → MAPE < 9%\n'
    '→ 输出95%置信区间', font_size=13, color=C_DARK)

# 异常检测部分
add_text_box(slide, 7, 1.4, 5.5, 3,
    '异常订单检测 — Isolation Forest\n\n'
    '核心直觉：异常点稀疏 → 随机切割 → 几下就隔离出来\n'
    '就像在白羊群中找黑羊比找另一只白羊更容易\n\n'
    '5维检测特征：\n'
    '• 订单金额\n'
    '• 商品种类数\n'
    '• 商品总数量\n'
    '• 平均单价\n'
    '• 折扣金额\n\n'
    '双层检测架构：\n'
    'IQR（单维初筛）→ Isolation Forest（多维复核）', font_size=13, color=C_DARK)

# 对比表
add_text_box(slide, 0.8, 4.8, 11, 0.5, '为什么不用LSTM/Prophet/ARIMA？', font_size=15, color=C_PRIMARY, bold=True)
add_text_box(slide, 0.8, 5.4, 11, 1.5,
    'LSTM：需要大数据量训练，90天不够 | 可解释性差 — "为什么预测这个值？"说不清楚\n'
    'Prophet：本质是调包，不能体现建模能力 | 依赖C++编译环境，部署容易出问题\n'
    'ARIMA：对节假日和特殊事件处理弱，需要人工哑变量 | 数据量小时不稳定\n'
    '→ 随机森林 + 特征工程 = 展示了我把时序转化为ML问题的通用方法论', font_size=13, color=C_DARK)

# ============================================================
# Slide 7: 报告 & 闭环
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_text_box(slide, 0.8, 0.5, 11, 0.6, '从数据到决策 — 分析闭环', font_size=28, color=C_DARK, bold=True)

# 闭环图
steps = ['上传\nCSV', '清洗\nETL', '分析\n建模', '可视化\nPlotly', '报告\n诊断']
for i, s in enumerate(steps):
    x = 1 + i * 2.3
    add_card(slide, x, 1.5, 1.8, 1.2, '', s, '', C_PRIMARY if i < 3 else C_ACCENT)
    if i < len(steps) - 1:
        add_text_box(slide, x + 1.85, 1.9, 0.3, 0.3, '→', font_size=24, color=C_GRAY, align=PP_ALIGN.CENTER)

# 报告示例
add_shape(slide, 0.8, 3.2, 11.7, 3.5, C_BG)
add_text_box(slide, 1.2, 3.3, 5, 3.2,
    '自动生成的诊断报告包含：\n\n'
    ' 经营概览\n'
    ' 总营收、客单价、退款率\n'
    ' 今日vs昨日对比\n'
    ' 统计五数（均值/中位数/偏度/峰度）\n\n'
    ' 异常提醒\n'
    ' 异常订单数量\n'
    ' 退款率偏高/正常判断\n'
    ' 日环比异常波动提示', font_size=13, color=C_DARK)

add_text_box(slide, 6.5, 3.3, 5.5, 3.2,
    ' 商品分析\n'
    ' 畅销Top5 & 滞销关注\n'
    ' 最佳套餐搭配建议\n\n'
    ' 用户洞察\n'
    ' 用户分层分布\n'
    ' 需召回的高价值客户数\n'
    ' 新客户引导建议\n\n'
    ' 预测 & 建议\n'
    ' 未来7天日均预测\n'
    ' 趋势判断（上升/下降/平稳）\n'
    ' 可执行的经营建议', font_size=13, color=C_DARK)

add_text_box(slide, 0.8, 6.3, 11, 0.8,
    '核心观点：数据分析的最终价值不在于出图表，而在于输出可执行的决策。这个项目完成了从原始CSV到自然语言经营建议的完整闭环。',
    font_size=14, color=C_PRIMARY, bold=True)

# ============================================================
# Slide 8: 项目亮点 & 技术总结
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_text_box(slide, 0.8, 0.5, 11, 0.6, '项目亮点 & 技术总结', font_size=28, color=C_WHITE, bold=True)

highlights = [
    ('全链路闭环', '从原始CSV上传到自然语言经营建议，\n覆盖6大分析页面+8个核心模块'),
    ('算法多样性', '数据挖掘(Apriori)+机器学习(K-Means/RF)\n+异常检测(Isolation Forest) 三类共6种算法'),
    ('特征工程', 'RFM三维特征构造、17维时间特征工程、\n时段特征衍生、品类偏好特征'),
    ('工程落地', '模块化架构(src/与pages/分离)、\nStreamlit Cloud公网部署、代码可维护'),
    ('数据思维', '多算法协同验证(RFM+KMeans交叉验证)\n双层异常检测(IQR初筛+IF多维复核)'),
    ('闭环设计', '从数据分析到可执行决策的最后一公里\n→ 一键生成自然语言经营建议报告'),
]

for i, (title, desc) in enumerate(highlights):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 1.5 + row * 2.8
    add_card(slide, x, y, 3.8, 2.3, '', title, desc, C_PRIMARY if row == 0 else C_ACCENT)

add_text_box(slide, 0.8, 6.8, 11, 0.4, 'restaurant-analytics-8fgygzvdk9b9edpxe4rpdc.streamlit.app',
           font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

# ===== 保存 =====
out = os.path.join(os.path.expanduser('~/Desktop'), '餐饮数据分析系统_项目展示.pptx')
prs.save(out)
print(f'Saved: {out}')
